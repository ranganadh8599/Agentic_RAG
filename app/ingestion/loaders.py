# Agentic RAG - document loaders.
# PDF: extracts text per page AND summarizes embedded images with a vision LLM
#      (pypdf text + page.images with an OOM guard, plus vision summarization
#      stored as sections).
# Images: vision-LLM extraction/description (works with any vision model).
# Text/Markdown: direct read.

import base64
import io
import os

from PIL import Image

from app.core.config import settings
from app.llm.client import chat_text


# ---------------------------------------------------------------------------
# Loaders return a list of sections: [{"text": str, "metadata": {...}}, ...]
# ---------------------------------------------------------------------------

def load_pdf(path, image_callback=None, max_pages: int = 0):
    from pypdf import PdfReader

    reader = PdfReader(path)
    sections = []
    total = len(reader.pages)
    if max_pages > 0:
        total = min(total, max_pages)

    # Per-DOCUMENT cap on vision calls (not per-page). Text extraction is
    # NEVER capped — only image summarization is.
    image_count = 0
    for page_no in range(total):
        page = reader.pages[page_no]
        text = (page.extract_text() or "").strip()
        if text:
            sections.append({"text": text, "metadata": {"page": page_no + 1, "kind": "text"}})

        # Embedded images -> vision summarization (with an OOM guard).
        if image_count < settings.MAX_IMAGES_PER_PDF:
            images = getattr(page, "images", None) or []
            for idx, img in enumerate(images):
                if image_count >= settings.MAX_IMAGES_PER_PDF:
                    break
                try:
                    data = img.data
                except Exception:
                    data = None
                if data:
                    result = _vision_summarize_image_bytes(data)
                    if result:
                        summary, img_bytes = result
                        image_count += 1
                        sections.append({
                            "text": summary,
                            "metadata": {"page": page_no + 1, "kind": "image", "image_idx": idx},
                            "image": {"data": img_bytes, "mime": "image/jpeg", "page": page_no + 1},
                        })
                        if image_callback:
                            image_callback(page_no + 1, image_count)
    return sections


def load_image(path):
    with open(path, "rb") as f:
        data = f.read()
    result = _vision_summarize_image_bytes(data)
    if not result:
        return []
    summary, img_bytes = result
    return [{"text": summary, "metadata": {"kind": "image"},
             "image": {"data": img_bytes, "mime": "image/jpeg", "page": None}}]


def load_text(path):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    return [{"text": content, "metadata": {"kind": "text"}}]


def load_docx(path):
    """Extract text from a Word .docx (paragraphs + tables)."""
    from docx import Document

    doc = Document(path)
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    text = "\n".join(parts)
    if not text.strip():
        return []
    return [{"text": text, "metadata": {"kind": "document"}}]


def load_xlsx(path):
    """Extract text from an Excel .xlsx workbook (one section per sheet)."""
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    sections = []
    try:
        for ws in wb.worksheets:
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) for v in row if v is not None]
                if vals:
                    rows_text.append(" | ".join(vals))
            if rows_text:
                header = f"[Sheet: {ws.title}]"
                sections.append({
                    "text": header + "\n" + "\n".join(rows_text),
                    "metadata": {"kind": "spreadsheet", "sheet": ws.title},
                })
    finally:
        wb.close()
    return sections


def load_pptx(path):
    """Extract text from a PowerPoint .pptx (one section per slide)."""
    from pptx import Presentation

    prs = Presentation(path)
    sections = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        parts.append(t)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        if parts:
            sections.append({"text": "\n".join(parts),
                             "metadata": {"kind": "slides", "slide": idx}})
    return sections


# ---------------------------------------------------------------------------
# Vision helper
# ---------------------------------------------------------------------------

def _vision_summarize_image_bytes(data: bytes):
    """Send an image to the vision model. Returns (summary_text, image_bytes)
    or None. The image bytes (resized JPEG) are kept so we can display the
    original alongside the answer."""
    # Resize huge images down so they fit API limits.
    try:
        img = Image.open(io.BytesIO(data))
        img.thumbnail((settings.MAX_IMAGE_DIM, settings.MAX_IMAGE_DIM))
        buf = io.BytesIO()
        img.convert("RGB").save(buf, format="JPEG", quality=settings.IMAGE_JPEG_QUALITY)
        data = buf.getvalue()
    except Exception:
        pass

    b64 = base64.b64encode(data).decode("ascii")
    # Gemini best practice: put the text prompt BEFORE the image (also removes
    # litellm's 'no text in user content' warning).
    instruction = (
        "Extract ALL readable text from this image and describe its visual "
        "content precisely and concisely. Output a clean text block suitable "
        "for search indexing."
    )
    messages = [
        {"role": "system", "content": "You are a precise vision assistant for document indexing."},
        {"role": "user", "content": [
            {"type": "text", "text": instruction},
            {"type": "image_url",
             "image_url": {"url": f"data:image/jpeg;base64,{b64}"}},
        ]},
    ]
    try:
        summary = chat_text(messages, model=settings.VISION_MODEL,
                            max_tokens=settings.VISION_SUMMARY_MAX_TOKENS)
    except Exception as exc:  # noqa: BLE001
        print(f"[vision] failed for image: {exc}")
        return None
    return (summary, data)
