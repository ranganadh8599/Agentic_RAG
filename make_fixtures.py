# Generate realistic test fixtures for the agentic-rag demo:
#   fixtures/report.pdf   - a 2-page PDF with facts about a fictional company
#   fixtures/notes.txt    - a plain-text document
#   fixtures/chart.png    - a simple generated image
# Run:  python make_fixtures.py

import os

from PIL import Image, ImageDraw
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

FIXTURES = os.path.join(os.path.dirname(__file__), "fixtures")
os.makedirs(FIXTURES, exist_ok=True)

REPORT_TEXT = (
    "Acme Analytics Annual Report 2024\n"
    "Executive Summary\n"
    "Acme Analytics reported total revenue of $2,400,000 in 2024, an increase of "
    "35% compared to 2023. The company added 1,200 new enterprise customers and "
    "expanded into three new markets: Germany, Japan, and Brazil.\n"
    "\n"
    "Products\n"
    "The flagship product, AcmeInsight, is an AI-powered business intelligence "
    "platform. It processes over 500 million data points per day. The mobile app "
    "was downloaded 800,000 times in 2024.\n"
    "\n"
    "Team\n"
    "Acme Analytics employs 240 people across offices in San Francisco, Berlin, "
    "and Singapore. The CTO, Dr. Maria Chen, leads a team of 60 engineers.\n"
    "\n"
    "Sustainability\n"
    "The company committed to achieving net-zero carbon emissions by 2030 and "
    "reduced energy consumption by 18% in 2024.\n"
)

NOTES_TEXT = (
    "Project Brainstorm Notes\n"
    "Idea: a multi-agent RAG system that reads PDFs and images and stores "
    "everything in PostgreSQL. It should support OpenAI, Gemini, and Claude. "
    "Key features: semantic caching, hybrid retrieval, grounded citations.\n"
    "\n"
    "TODOs:\n"
    "- Try pgvector HNSW indexes for fast search\n"
    "- Add a vision model to summarize embedded images\n"
    "- Compare char vs BPE tokenizers from the gpt-from-scratch project\n"
)


def make_pdf():
    path = os.path.join(FIXTURES, "report.pdf")
    c = canvas.Canvas(path, pagesize=letter)
    width, height = letter
    # Page 1
    for i, line in enumerate(REPORT_TEXT.split("\n")):
        c.drawString(72, height - 72 - i * 16, line)
    c.showPage()
    # Page 2
    c.drawString(72, height - 72, "Financial Details")
    c.drawString(72, height - 96, "Revenue breakdown by product:")
    c.drawString(72, height - 120, "- AcmeInsight subscriptions: $1,600,000")
    c.drawString(72, height - 144, "- Consulting services: $600,000")
    c.drawString(72, height - 168, "- Licenses: $200,000")
    c.save()
    print("wrote", path)


def make_image():
    path = os.path.join(FIXTURES, "chart.png")
    img = Image.new("RGB", (640, 400), "white")
    d = ImageDraw.Draw(img)
    d.rectangle([40, 40, 600, 360], outline="black")
    d.rectangle([80, 220, 160, 340], fill="steelblue")
    d.rectangle([200, 160, 280, 340], fill="steelblue")
    d.rectangle([320, 100, 400, 340], fill="steelblue")
    d.text((90, 345), "Q1", fill="black")
    d.text((215, 345), "Q2", fill="black")
    d.text((335, 345), "Q3", fill="black")
    d.text((60, 60), "Acme Analytics Quarterly Revenue", fill="black")
    d.text((80, 200), "2", fill="black")
    d.text((200, 140), "4", fill="black")
    d.text((320, 80), "6", fill="black")
    img.save(path)
    print("wrote", path)


def make_text():
    path = os.path.join(FIXTURES, "notes.txt")
    with open(path, "w", encoding="utf-8") as f:
        f.write(NOTES_TEXT)
    print("wrote", path)


def make_office():
    """Create .docx, .xlsx and .pptx test files."""
    # DOCX
    from docx import Document
    d = Document()
    d.add_heading("Quarterly Business Review", level=1)
    d.add_paragraph("The team grew revenue by 22% this quarter.")
    d.add_paragraph("Key initiative: launch of the mobile app.")
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "Metric"
    t.cell(0, 1).text = "Value"
    t.cell(1, 0).text = "Revenue"
    t.cell(1, 1).text = "$1.2M"
    d.save(os.path.join(FIXTURES, "business.docx"))

    # XLSX
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Region", "Q1", "Q2"])
    ws.append(["North", 100, 120])
    ws.append(["South", 90, 130])
    wb.save(os.path.join(FIXTURES, "sales.xlsx"))

    # PPTX
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Product Roadmap"
    slide.placeholders[1].text = "Phase 1: AI features. Phase 2: Mobile."
    prs.save(os.path.join(FIXTURES, "roadmap.pptx"))

    for f in ("business.docx", "sales.xlsx", "roadmap.pptx"):
        print("wrote", os.path.join(FIXTURES, f))


def make_pdf_with_image():
    """Create a PDF that embeds chart.png so we can test embedded-image display."""
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    path = os.path.join(FIXTURES, "report_with_image.pdf")
    img_path = os.path.join(FIXTURES, "chart.png")
    c = canvas.Canvas(path, pagesize=letter)
    width, height = letter
    c.drawString(72, height - 72, "Quarterly Revenue Summary")
    c.drawString(72, height - 96, "The chart below shows revenue growth by quarter.")
    c.drawImage(img_path, 72, height - 420, width=300, height=200,
                preserveAspectRatio=True, mask="auto")
    c.showPage()
    c.save()
    print("wrote", path)


if __name__ == "__main__":
    make_pdf()
    make_image()
    make_text()
    make_office()
    make_pdf_with_image()
