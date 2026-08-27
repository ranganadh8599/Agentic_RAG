"""Unit tests for the document loaders.

Sample files are generated in-memory with the same libraries the loaders use
(docx / openpyxl / python-pptx / reportlab / PIL), so no binary fixtures need
to be committed. VISION_MODEL is mock in the test env, so load_image runs
offline.
"""
import os
import tempfile

import pytest
from PIL import Image

from app.ingestion.loaders import (load_docx, load_image, load_pdf, load_pptx,
                                   load_text, load_xlsx)


def _tmp(suffix):
    return tempfile.mktemp(suffix=suffix)


def test_load_text():
    path = _tmp(".txt")
    try:
        with open(path, "w", encoding="utf-8") as f:
            f.write("hello world\nsecond line")
        sections = load_text(path)
        assert sections[0]["text"] == "hello world\nsecond line"
        assert sections[0]["metadata"]["kind"] == "text"
    finally:
        os.unlink(path)


def test_load_text_missing_file_raises():
    with pytest.raises(FileNotFoundError):
        load_text(_tmp(".txt"))


def test_load_docx():
    from docx import Document
    path = _tmp(".docx")
    doc = Document()
    doc.add_paragraph("Paragraph one.")
    doc.add_paragraph("Paragraph two.")
    doc.save(path)
    try:
        sections = load_docx(path)
        assert "Paragraph one." in sections[0]["text"]
        assert "Paragraph two." in sections[0]["text"]
    finally:
        os.unlink(path)


def test_load_xlsx():
    from openpyxl import Workbook
    path = _tmp(".xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Value"])
    ws.append(["Alpha", 42])
    wb.save(path)
    try:
        sections = load_xlsx(path)
        assert "[Sheet: Sheet1]" in sections[0]["text"]
        assert "Alpha" in sections[0]["text"]
        assert "42" in sections[0]["text"]
    finally:
        os.unlink(path)


def test_load_pptx():
    from pptx import Presentation
    path = _tmp(".pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello Slide"
    prs.save(path)
    try:
        sections = load_pptx(path)
        assert any("Hello Slide" in s["text"] for s in sections)
    finally:
        os.unlink(path)


def test_load_pdf():
    from reportlab.pdfgen import canvas
    path = _tmp(".pdf")
    c = canvas.Canvas(path)
    c.drawString(72, 720, "Extract this PDF text.")
    c.save()
    try:
        sections = load_pdf(path, max_pages=0)
        texts = " ".join(s["text"] for s in sections)
        assert "Extract this PDF text." in texts
    finally:
        os.unlink(path)


def test_load_image_uses_vision_model():
    path = _tmp(".png")
    Image.new("RGB", (16, 16), color=(255, 0, 0)).save(path)
    try:
        sections = load_image(path)
        # VISION_MODEL=mock → the mock chat reply is returned as the summary.
        assert sections and sections[0]["text"]
        assert sections[0]["metadata"]["kind"] == "image"
    finally:
        os.unlink(path)
